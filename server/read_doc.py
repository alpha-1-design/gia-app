#!/usr/bin/env python3
"""Universal document reader for GIA.
Reads any input document format via stdin (as base64) or command-line path argument.
Outputs plain text to stdout.

Supported formats: DOCX, XLSX, PPTX, ODT, ODS, EPUB, HTML, Markdown, RTF, CSV, JSON, XML, YAML, plain text
"""
import sys
import os
import json
import base64
import struct
import zipfile
import xml.etree.ElementTree as ET
import io
import csv
import re
from pathlib import Path

def die(msg):
    print(json.dumps({"error": msg}), file=sys.stderr)
    sys.exit(1)

def safe_text(elem, sep=" "):
    """Extract text from an XML element recursively."""
    parts = []
    if elem.text:
        parts.append(elem.text.strip())
    for child in elem:
        parts.append(safe_text(child, sep))
        if child.tail:
            parts.append(child.tail.strip())
    return sep.join(filter(None, parts))

def parse_docx(path):
    """Microsoft Word (.docx)"""
    from docx import Document
    doc = Document(path)
    lines = []
    for para in doc.paragraphs:
        lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))
    return "\n".join(lines)

def parse_xlsx(path):
    """Microsoft Excel (.xlsx)"""
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"=== Sheet: {sheet_name} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) if v is not None else "" for v in row]
            lines.append("\t".join(vals))
    return "\n".join(lines)

def parse_pptx(path):
    """Microsoft PowerPoint (.pptx)"""
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append(" | ".join(cells))
    return "\n".join(lines)

def parse_odt(path):
    """OpenDocument Text (.odt)"""
    import odf.opendocument as od
    from odf.text import P
    doc = od.load(path)
    lines = []
    for elem in doc.getElementsByType(P):
        text = safe_text(elem)
        if text:
            lines.append(text)
    return "\n".join(lines)

def parse_ods(path):
    """OpenDocument Spreadsheet (.ods)"""
    import odf.opendocument as od
    from odf.table import Table, TableRow, TableCell
    from odf.text import P as TextP
    doc = od.load(path)
    lines = []
    for table in doc.getElementsByType(Table):
        name = table.getAttribute("name") or "Sheet"
        lines.append(f"=== Sheet: {name} ===")
        for row in table.getElementsByType(TableRow):
            cells = []
            for cell in row.getElementsByType(TableCell):
                texts = []
                for p in cell.getElementsByType(TextP):
                    t = safe_text(p)
                    if t:
                        texts.append(t)
                cells.append(" ".join(texts))
            lines.append("\t".join(cells))
    return "\n".join(lines)

def parse_epub(path):
    """EPUB ebook"""
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    book = epub.read_epub(path)
    lines = []
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), "html.parser")
            text = soup.get_text(separator="\n").strip()
            if text:
                lines.append(text)
    return "\n".join(lines)

def parse_html(path):
    """HTML file"""
    from bs4 import BeautifulSoup
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    return soup.get_text(separator="\n").strip()

def parse_markdown(path):
    """Markdown file"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_rtf(path):
    """RTF - rough text extraction"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    text = re.sub(r"\\[a-z]+[-0-9]*", " ", text)
    text = re.sub(r"\{|\}", " ", text)
    text = re.sub(r"\\'[0-9a-f]{2}", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text

def parse_csv(path):
    """CSV file"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        lines = []
        for row in reader:
            lines.append(", ".join(row))
        return "\n".join(lines)

def parse_json(path):
    """JSON file - pretty printed"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)

def parse_xml(path):
    """XML file"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_yaml(path):
    """YAML file"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()

def parse_txt(path):
    """Plain text file"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


EXT_MAP = {
    ".docx": ("docx", "Microsoft Word Document", parse_docx),
    ".xlsx": ("xlsx", "Microsoft Excel Spreadsheet", parse_xlsx),
    ".pptx": ("pptx", "Microsoft PowerPoint Presentation", parse_pptx),
    ".odt": ("odt", "OpenDocument Text", parse_odt),
    ".ods": ("ods", "OpenDocument Spreadsheet", parse_ods),
    ".epub": ("epub", "EPUB eBook", parse_epub),
    ".html": ("html", "HTML Document", parse_html),
    ".htm": ("html", "HTML Document", parse_html),
    ".md": ("md", "Markdown", parse_markdown),
    ".markdown": ("md", "Markdown", parse_markdown),
    ".rtf": ("rtf", "Rich Text Format", parse_rtf),
    ".csv": ("csv", "CSV Spreadsheet", parse_csv),
    ".json": ("json", "JSON Data", parse_json),
    ".xml": ("xml", "XML Document", parse_xml),
    ".yaml": ("yaml", "YAML Data", parse_yaml),
    ".yml": ("yaml", "YAML Data", parse_yaml),
    ".txt": ("txt", "Plain Text", parse_txt),
    ".log": ("txt", "Log File", parse_txt),
    ".env": ("txt", "Environment File", parse_txt),
    ".cfg": ("txt", "Configuration File", parse_txt),
    ".ini": ("txt", "INI Configuration", parse_txt),
    ".conf": ("txt", "Configuration File", parse_txt),
    ".py": ("txt", "Python Script", parse_txt),
    ".js": ("txt", "JavaScript", parse_txt),
    ".ts": ("txt", "TypeScript", parse_txt),
    ".jsx": ("txt", "JSX", parse_txt),
    ".tsx": ("txt", "TSX", parse_txt),
    ".css": ("txt", "CSS", parse_txt),
    ".scss": ("txt", "SCSS", parse_txt),
    ".sh": ("txt", "Shell Script", parse_txt),
    ".bat": ("txt", "Batch Script", parse_txt),
    ".ps1": ("txt", "PowerShell", parse_txt),
    ".sql": ("txt", "SQL Script", parse_txt),
    ".java": ("txt", "Java Source", parse_txt),
    ".c": ("txt", "C Source", parse_txt),
    ".cpp": ("txt", "C++ Source", parse_txt),
    ".h": ("txt", "C Header", parse_txt),
    ".hpp": ("txt", "C++ Header", parse_txt),
    ".go": ("txt", "Go Source", parse_txt),
    ".rs": ("txt", "Rust Source", parse_txt),
    ".rb": ("txt", "Ruby Source", parse_txt),
    ".php": ("txt", "PHP Source", parse_txt),
    ".swift": ("txt", "Swift Source", parse_txt),
    ".kt": ("txt", "Kotlin Source", parse_txt),
    ".dart": ("txt", "Dart Source", parse_txt),
    ".lua": ("txt", "Lua Script", parse_txt),
    ".r": ("txt", "R Script", parse_txt),
    ".toml": ("txt", "TOML Config", parse_txt),
    ".dockerfile": ("txt", "Dockerfile", parse_txt),
    ".gitignore": ("txt", "Git Ignore", parse_txt),
    ".svg": ("txt", "SVG Image", parse_txt),
    ".tex": ("txt", "LaTeX Document", parse_txt),
}

def main():
    # Determine input path
    if len(sys.argv) > 1:
        path = sys.argv[1]
        if not os.path.exists(path):
            die(f"File not found: {path}")
        ext = os.path.splitext(path)[1].lower()
    else:
        # Read from stdin: either base64 data or JSON with {path, data}
        raw = sys.stdin.buffer.read()
        try:
            decoded = raw.decode("utf-8").strip()
        except UnicodeDecodeError:
            decoded = ""
        
        # Check if stdin is a JSON request with path
        if decoded.startswith("{"):
            try:
                req = json.loads(decoded)
            except json.JSONDecodeError:
                req = {}
            filepath = req.get("path", "")
            file_data = req.get("data", "")
            if filepath:
                path = filepath
                ext = os.path.splitext(path)[1].lower()
                # Write data to temp file if provided
                if file_data:
                    temp = f"/tmp/gia_doc_{os.getpid()}{ext}"
                    with open(temp, "wb") as f:
                        f.write(base64.b64decode(file_data))
                    path = temp
                if not os.path.exists(path):
                    die(f"File not found: {path}")
            else:
                # Try to interpret stdin directly as path
                path = decoded
                if not os.path.exists(path):
                    die(f"File not found and no path/data in JSON: {path}")
                ext = os.path.splitext(path)[1].lower()
        else:
            path = decoded
            if not os.path.exists(path):
                die(f"File not found: {path}")
            ext = os.path.splitext(path)[1].lower()
    
    fmt = EXT_MAP.get(ext)
    if not fmt:
        # Try magic bytes
        with open(path, "rb") as f:
            magic = f.read(8)
        # ZIP-based formats (docx/xlsx/pptx/epub/odt/ods)
        if magic[:4] == b"PK\x03\x04":
            try:
                with zipfile.ZipFile(path) as z:
                    names = z.namelist()
                    if "word/document.xml" in names:
                        fmt = ("docx", "Microsoft Word Document", parse_docx)
                    elif "xl/workbook.xml" in names or "xl/sharedStrings.xml" in names:
                        fmt = ("xlsx", "Microsoft Excel Spreadsheet", parse_xlsx)
                    elif "ppt/presentation.xml" in names:
                        fmt = ("pptx", "Microsoft PowerPoint Presentation", parse_pptx)
                    elif "META-INF/container.xml" in names:
                        # Could be ODT or ODS
                        try:
                            cf = z.read("META-INF/container.xml")
                            if b"office:document" in cf or b"application/vnd.oasis" in cf:
                                # Read rootfile
                                root = ET.fromstring(cf)
                            fmt = ("odf", "OpenDocument", parse_odt)
                        except:
                            fmt = ("epub", "EPUB eBook", parse_epub)
                    elif "mimetype" in names:
                        mt = z.read("mimetype").decode("utf-8").strip()
                        if "epub" in mt:
                            fmt = ("epub", "EPUB eBook", parse_epub)
                        else:
                            fmt = ("odf", "OpenDocument", parse_odt)
                    else:
                        fmt = ("zip", "ZIP Archive", None)
            except:
                fmt = ("zip", "ZIP Archive", None)
        elif magic[:4] == b"%PDF":
            fmt = ("pdf", "PDF Document", None)
        elif magic[:4] == b"<!DO" or magic[:4] == b"<htm" or magic[:4].lower() == b"<htm":
            fmt = ("html", "HTML Document", parse_html)
        elif magic[:2] == b"\\v" or magic[:2] == b"{\\rt":
            fmt = ("rtf", "Rich Text Format", None)
        elif magic[:2] == b"\xff\xfe" or magic[:2] == b"\xfe\xff":
            fmt = ("txt", "UTF-16 Text", parse_txt)
        else:
            # Treat as plain text
            fmt = ("txt", "Plain Text", parse_txt)
    
    fmt_id, fmt_name, parser = fmt
    
    if parser is None:
        die(f"Reading .{fmt_id} ({fmt_name}) files is not supported yet. Try converting to a supported format.")
    
    try:
        text = parser(path)
        result = {
            "format": fmt_id,
            "format_name": fmt_name,
            "content": text,
            "size": len(text),
        }
        print(json.dumps(result, ensure_ascii=False))
    except ImportError as e:
        missing_pkg = str(e).split("'")[1] if "'" in str(e) else str(e)
        die(f"Missing Python package for {fmt_name} parsing: {missing_pkg}. Install with: pip install {missing_pkg}")
    except Exception as e:
        die(f"Error parsing {fmt_name}: {e}")

if __name__ == "__main__":
    main()
